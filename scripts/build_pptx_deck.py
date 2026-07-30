import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BG_DARK = RGBColor(15, 23, 42)
    TEXT_DARK = RGBColor(255, 255, 255)
    CARD_BG = RGBColor(30, 41, 59)
    ACCENT_BLUE = RGBColor(0, 117, 222)
    ACCENT_GREEN = RGBColor(16, 185, 129)
    TEXT_MUTED = RGBColor(148, 163, 184)

    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category_text="BUZZ AGENT WORKSPACE"):
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE

        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(26)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_DARK

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

    # --- SLIDE 1: Title Slide ---
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    shape = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = ACCENT_BLUE
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "BUZZ AGENT WORKSPACE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "\nAutonomous Software, Marketing & Job Hunter Fleets"
    p2.font.size = Pt(30)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_DARK
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "\nDecentralized Nostr Protocol • Hermes Agent OS on Hostinger • Human-in-the-Loop Assist"
    p3.font.size = Pt(16)
    p3.font.color.rgb = TEXT_MUTED
    p3.alignment = PP_ALIGN.CENTER

    # --- SLIDE 2: Surface Architecture ---
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "Surface Roles & System Division of Labor")

    surfaces = [
        ("Buzz App (buzz.xyz)", "Human Workspace & Chat GUI", "Nostr WebSocket client. Real-time channel timeline (#engineering, #marketing-ops, #job-hunter), 1-click approvals, and screenshot previews.", ACCENT_BLUE),
        ("Hermes Agent OS", "Remote Execution Engine (Hostinger)", "Runs headless on devserver. Subscribes to Nostr relay via buzz-acp, executes skills, manages memory & local Ollama inference.", ACCENT_GREEN),
        ("AntiGravity / IDE", "Code Review & Navigation Pane", "Human pair-programming workspace for codebase inspection, architectural design, quality gates, and git repository control.", ACCENT_BLUE)
    ]

    for i, (name, role, desc, color) in enumerate(surfaces):
        left = Inches(0.8 + i * 3.9)
        top = Inches(1.8)
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.7), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

        p_role = tf.add_paragraph()
        p_role.text = f"\n{role}"
        p_role.font.size = Pt(13)
        p_role.font.bold = True
        p_role.font.color.rgb = color

        p_desc = tf.add_paragraph()
        p_desc.text = f"\n{desc}"
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_MUTED

    # --- SLIDE 3: The 3 Agent Fleets ---
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "The 3 Autonomous Agent Fleets")

    fleets = [
        ("💻 Software Factory", "@SystemArchitect • @FeatureDeveloper\n@QAGatekeeper (P-3) • @DevOpsRelease", "TDD code generation, feature branches, four-eyes code review gates, and staging deployments."),
        ("📈 Marketing Factory", "@GrowthAnalyst • @SEOIntelAgent\n@CreativeCopywriter • @AdDeployer", "SEMrush/SE Ranking scans, Notion-Warm landing pages, proof-based ad copy, and Meta/Google Ads drafting."),
        ("🍂 AutHarvest Job Hunter", "@JobScanner • @MatchScorer\n@CVTailor • @ApplyAssistant", "Automated job discovery, £80K+/£600+day scoring, fact-grounded CV tailoring, and 1-click application assist.")
    ]

    for i, (name, agents, desc) in enumerate(fleets):
        left = Inches(0.8 + i * 3.9)
        top = Inches(1.8)
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.7), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = ACCENT_BLUE
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

        p_ag = tf.add_paragraph()
        p_ag.text = f"\nAgents:\n{agents}"
        p_ag.font.size = Pt(12)
        p_ag.font.bold = True
        p_ag.font.color.rgb = ACCENT_GREEN

        p_desc = tf.add_paragraph()
        p_desc.text = f"\nCapabilities:\n{desc}"
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_MUTED

    # --- SLIDE 4: AutHarvest Workflow ---
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "AutHarvest: Human-in-the-Loop Application Staging")

    steps = [
        ("1. Discovery & Scoring", "07:00 UTC Cron runs @JobScanner (LinkedIn/Indeed) and @MatchScorer (>=4.0 threshold)."),
        ("2. Fact-Grounded CV", "@CVTailor builds tailored .docx CV & Cover Letter with provenance check against ruvector.db."),
        ("3. NIP-44 Encrypted Card", "Private encrypted Nostr event sent to #job-hunter with Pocket vs Launch choices."),
        ("4. 5-Sec Application Assist", "Clicking Launch opens browser to job URL with fields pre-filled for 1-click human submit.")
    ]

    for i, (title, desc) in enumerate(steps):
        top = Inches(1.8 + i * 1.3)
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(1.1))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = ACCENT_BLUE
        card.line.width = Pt(1)

        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_DARK

    # --- SLIDE 5: Security & Mama Obsidian ---
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "Security Safeguards & Mama Obsidian Vault Sync")

    rules = [
        ("🔒 NIP-44 Pairwise Encryption", "All candidate PII, job summaries, and CV links are encrypted before relay broadcast."),
        ("🛡️ P-3 Four-Eyes Review Gate", "Non-Anthropic models (Qwen 3.7 / Gemma 3) independently audit code and campaigns before release."),
        ("🧠 Mama Obsidian Vault Sync", "Every agent output logs structured CommonMark notes to /second-brain/ for local Ollama RAG."),
        ("⚡ 100% Compute Portability", "Agents run location-agnostically on Hostinger devserver, GCP Nano VMs, or cloud GPU rigs.")
    ]

    for i, (title, desc) in enumerate(rules):
        col = i % 2
        row = i // 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.8 + row * 2.5)
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.6), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = ACCENT_BLUE
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

        p_desc = tf.add_paragraph()
        p_desc.text = f"\n{desc}"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MUTED

    os.makedirs("/Users/arajiv/buzz-implementation-plan/outputs", exist_ok=True)
    out_path = "/Users/arajiv/buzz-implementation-plan/outputs/Buzz_Architecture_Deck.pptx"
    prs.save(out_path)
    print(f"Presentation saved successfully to {out_path}")

if __name__ == "__main__":
    create_deck()
